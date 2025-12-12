"""
PDF Toolkit Backend - Flask Application
A comprehensive PDF processing API with multiple conversion and manipulation features.
"""

import os
import io
import uuid
import shutil
import tempfile
import zipfile
from pathlib import Path
from typing import List, Optional
from datetime import datetime

from flask import Flask, request, jsonify, send_file, Response, send_from_directory
from flask_cors import CORS
from werkzeug.utils import secure_filename

import pypdf
from pypdf import PdfReader, PdfWriter
from pdf2docx import Converter as PDF2DocxConverter
from pdf2image import convert_from_path
from PIL import Image, ImageDraw, ImageFont
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from reportlab.lib.colors import Color
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import pikepdf
import requests as http_requests
from bs4 import BeautifulSoup
import pdfkit
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import openpyxl
from openpyxl import Workbook
import fitz  # PyMuPDF
import subprocess

app = Flask(__name__, static_folder=None)
CORS(app)

# Directories
BASE_DIR = Path(__file__).parent
UPLOAD_DIR = BASE_DIR / "uploads"
OUTPUT_DIR = BASE_DIR / "outputs"
FRONTEND_DIR = BASE_DIR / "frontend"
UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)


# ==================== SERVE FRONTEND ====================
@app.route("/")
def serve_index():
    """Serve the frontend index.html."""
    return send_from_directory(FRONTEND_DIR, "index.html")


@app.route("/<path:filename>")
def serve_static(filename):
    """Serve static files from the frontend directory."""
    return send_from_directory(FRONTEND_DIR, filename)


def generate_filename(prefix: str, extension: str) -> str:
    """Generate a unique filename with timestamp."""
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    unique_id = str(uuid.uuid4())[:8]
    return f"{prefix}_{timestamp}_{unique_id}.{extension}"


def cleanup_file(filepath: Path):
    """Remove a file after download."""
    try:
        if filepath.exists():
            filepath.unlink()
    except Exception:
        pass


def save_upload_file(upload_file) -> Path:
    """Save an uploaded file and return its path."""
    filename = secure_filename(upload_file.filename)
    file_path = UPLOAD_DIR / f"{uuid.uuid4()}_{filename}"
    upload_file.save(str(file_path))
    return file_path


# ==================== MERGE PDF ====================
@app.route("/api/merge", methods=["POST"])
def merge_pdfs():
    """Merge multiple PDF files into one."""
    files = request.files.getlist("files")
    
    if len(files) < 2:
        return jsonify({"detail": "At least 2 PDF files required"}), 400
    
    writer = PdfWriter()
    temp_files = []
    
    try:
        for file in files:
            temp_path = save_upload_file(file)
            temp_files.append(temp_path)
            reader = PdfReader(str(temp_path))
            for page in reader.pages:
                writer.add_page(page)
        
        output_filename = generate_filename("merged", "pdf")
        output_path = OUTPUT_DIR / output_filename
        with open(output_path, "wb") as f:
            writer.write(f)
        
        response = send_file(
            output_path,
            mimetype="application/pdf",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    finally:
        for temp_file in temp_files:
            cleanup_file(temp_file)


# ==================== SPLIT PDF ====================
@app.route("/api/split", methods=["POST"])
def split_pdf():
    """Split a PDF into individual pages or extract specific pages."""
    file = request.files.get("file")
    pages = request.form.get("pages", "all")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    temp_path = save_upload_file(file)
    
    try:
        reader = PdfReader(str(temp_path))
        total_pages = len(reader.pages)
        
        # Parse page selection
        if pages == "all":
            selected_pages = list(range(total_pages))
        else:
            selected_pages = []
            for part in pages.split(","):
                part = part.strip()
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    selected_pages.extend(range(start - 1, min(end, total_pages)))
                else:
                    page_num = int(part) - 1
                    if 0 <= page_num < total_pages:
                        selected_pages.append(page_num)
        
        # Create ZIP with split pages
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            for i, page_num in enumerate(selected_pages):
                writer = PdfWriter()
                writer.add_page(reader.pages[page_num])
                
                page_buffer = io.BytesIO()
                writer.write(page_buffer)
                page_buffer.seek(0)
                
                zip_file.writestr(f"page_{page_num + 1}.pdf", page_buffer.getvalue())
        
        zip_buffer.seek(0)
        return Response(
            zip_buffer.getvalue(),
            mimetype="application/zip",
            headers={"Content-Disposition": "attachment; filename=split_pages.zip"}
        )
    finally:
        cleanup_file(temp_path)


# ==================== PDF TO WORD ====================
@app.route("/api/pdf-to-word", methods=["POST"])
def pdf_to_word():
    """Convert PDF to DOCX format."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    temp_path = save_upload_file(file)
    output_filename = generate_filename("converted", "docx")
    output_path = OUTPUT_DIR / output_filename
    
    try:
        cv = PDF2DocxConverter(str(temp_path))
        cv.convert(str(output_path))
        cv.close()
        
        response = send_file(
            output_path,
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    finally:
        cleanup_file(temp_path)


# ==================== PDF TO POWERPOINT ====================
@app.route("/api/pdf-to-ppt", methods=["POST"])
def pdf_to_powerpoint():
    """Convert PDF to PPTX format (each page becomes a slide with image)."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    temp_path = save_upload_file(file)
    output_filename = generate_filename("converted", "pptx")
    output_path = OUTPUT_DIR / output_filename
    
    try:
        # Convert PDF pages to images
        images = convert_from_path(str(temp_path), dpi=200)
        
        # Create PowerPoint
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        for i, img in enumerate(images):
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Save image temporarily
            img_path = UPLOAD_DIR / f"temp_slide_{i}.png"
            img.save(str(img_path), "PNG")
            
            # Add image to slide
            slide.shapes.add_picture(
                str(img_path),
                Inches(0),
                Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )
            cleanup_file(img_path)
        
        prs.save(str(output_path))
        
        response = send_file(
            output_path,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    finally:
        cleanup_file(temp_path)


# ==================== PDF TO EXCEL ====================
@app.route("/api/pdf-to-excel", methods=["POST"])
def pdf_to_excel():
    """Extract tables from PDF to Excel format."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    temp_path = save_upload_file(file)
    output_filename = generate_filename("converted", "xlsx")
    output_path = OUTPUT_DIR / output_filename
    
    try:
        # Use PyMuPDF to extract text and tables
        doc = fitz.open(str(temp_path))
        wb = Workbook()
        ws = wb.active
        ws.title = "PDF Content"
        
        row_num = 1
        for page_num, page in enumerate(doc):
            # Add page header
            ws.cell(row=row_num, column=1, value=f"--- Page {page_num + 1} ---")
            row_num += 1
            
            # Extract tables
            tables = page.find_tables()
            if tables:
                for table in tables:
                    for table_row in table.extract():
                        for col_num, cell in enumerate(table_row, 1):
                            ws.cell(row=row_num, column=col_num, value=cell)
                        row_num += 1
                    row_num += 1  # Empty row between tables
            else:
                # If no tables, extract text
                text = page.get_text()
                for line in text.split('\n'):
                    if line.strip():
                        ws.cell(row=row_num, column=1, value=line.strip())
                        row_num += 1
            
            row_num += 1  # Empty row between pages
        
        doc.close()
        wb.save(str(output_path))
        
        response = send_file(
            output_path,
            mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    finally:
        cleanup_file(temp_path)


# ==================== WORD TO PDF ====================
@app.route("/api/word-to-pdf", methods=["POST"])
def word_to_pdf():
    """Convert DOCX to PDF format."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    temp_path = save_upload_file(file)
    output_filename = generate_filename("converted", "pdf")
    output_path = OUTPUT_DIR / output_filename
    
    try:
        # Try using LibreOffice for conversion
        result = subprocess.run([
            "libreoffice", "--headless", "--convert-to", "pdf",
            "--outdir", str(OUTPUT_DIR), str(temp_path)
        ], capture_output=True, timeout=120)
        
        # Find the converted file
        converted_name = temp_path.stem + ".pdf"
        converted_path = OUTPUT_DIR / converted_name
        
        if converted_path.exists():
            shutil.move(str(converted_path), str(output_path))
        else:
            return jsonify({"detail": "Conversion failed"}), 500
        
        response = send_file(
            output_path,
            mimetype="application/pdf",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    except subprocess.TimeoutExpired:
        return jsonify({"detail": "Conversion timed out"}), 500
    except FileNotFoundError:
        return jsonify({"detail": "LibreOffice not installed"}), 500
    finally:
        cleanup_file(temp_path)


# ==================== POWERPOINT TO PDF ====================
@app.route("/api/ppt-to-pdf", methods=["POST"])
def ppt_to_pdf():
    """Convert PPTX to PDF format."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    temp_path = save_upload_file(file)
    output_filename = generate_filename("converted", "pdf")
    output_path = OUTPUT_DIR / output_filename
    
    try:
        result = subprocess.run([
            "libreoffice", "--headless", "--convert-to", "pdf",
            "--outdir", str(OUTPUT_DIR), str(temp_path)
        ], capture_output=True, timeout=120)
        
        converted_name = temp_path.stem + ".pdf"
        converted_path = OUTPUT_DIR / converted_name
        
        if converted_path.exists():
            shutil.move(str(converted_path), str(output_path))
        else:
            return jsonify({"detail": "Conversion failed"}), 500
        
        response = send_file(
            output_path,
            mimetype="application/pdf",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    except subprocess.TimeoutExpired:
        return jsonify({"detail": "Conversion timed out"}), 500
    except FileNotFoundError:
        return jsonify({"detail": "LibreOffice not installed"}), 500
    finally:
        cleanup_file(temp_path)


# ==================== EXCEL TO PDF ====================
@app.route("/api/excel-to-pdf", methods=["POST"])
def excel_to_pdf():
    """Convert XLSX to PDF format."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    temp_path = save_upload_file(file)
    output_filename = generate_filename("converted", "pdf")
    output_path = OUTPUT_DIR / output_filename
    
    try:
        result = subprocess.run([
            "libreoffice", "--headless", "--convert-to", "pdf",
            "--outdir", str(OUTPUT_DIR), str(temp_path)
        ], capture_output=True, timeout=120)
        
        converted_name = temp_path.stem + ".pdf"
        converted_path = OUTPUT_DIR / converted_name
        
        if converted_path.exists():
            shutil.move(str(converted_path), str(output_path))
        else:
            return jsonify({"detail": "Conversion failed"}), 500
        
        response = send_file(
            output_path,
            mimetype="application/pdf",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    except subprocess.TimeoutExpired:
        return jsonify({"detail": "Conversion timed out"}), 500
    except FileNotFoundError:
        return jsonify({"detail": "LibreOffice not installed"}), 500
    finally:
        cleanup_file(temp_path)


# ==================== EDIT PDF ====================
@app.route("/api/edit-pdf", methods=["POST"])
def edit_pdf():
    """Add text to a PDF document."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    text = request.form.get("text", "")
    text_x = float(request.form.get("text_x", 100))
    text_y = float(request.form.get("text_y", 100))
    text_size = int(request.form.get("text_size", 12))
    text_color = request.form.get("text_color", "#000000")
    page_num = int(request.form.get("page_num", 1))
    
    temp_path = save_upload_file(file)
    output_filename = generate_filename("edited", "pdf")
    output_path = OUTPUT_DIR / output_filename
    
    try:
        doc = fitz.open(str(temp_path))
        
        if page_num < 1 or page_num > len(doc):
            return jsonify({"detail": "Invalid page number"}), 400
        
        page = doc[page_num - 1]
        
        # Parse color
        color_hex = text_color.lstrip("#")
        r = int(color_hex[0:2], 16) / 255
        g = int(color_hex[2:4], 16) / 255
        b = int(color_hex[4:6], 16) / 255
        
        # Insert text
        if text:
            text_point = fitz.Point(text_x, text_y)
            page.insert_text(
                text_point,
                text,
                fontsize=text_size,
                color=(r, g, b)
            )
        
        doc.save(str(output_path))
        doc.close()
        
        response = send_file(
            output_path,
            mimetype="application/pdf",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    finally:
        cleanup_file(temp_path)


# ==================== PDF TO JPG ====================
@app.route("/api/pdf-to-jpg", methods=["POST"])
def pdf_to_jpg():
    """Convert PDF pages to JPG images or extract embedded images."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    dpi = int(request.form.get("dpi", 200))
    extract_images = request.form.get("extract_images", "false").lower() in ("true", "1", "yes")
    
    temp_path = save_upload_file(file)
    
    try:
        zip_buffer = io.BytesIO()
        
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            if extract_images:
                # Extract embedded images
                doc = fitz.open(str(temp_path))
                img_count = 0
                
                for page_num, page in enumerate(doc):
                    image_list = page.get_images()
                    for img_index, img in enumerate(image_list):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        
                        # Convert to JPG if not already
                        pil_img = Image.open(io.BytesIO(image_bytes))
                        if pil_img.mode in ("RGBA", "P"):
                            pil_img = pil_img.convert("RGB")
                        
                        img_buffer = io.BytesIO()
                        pil_img.save(img_buffer, format="JPEG", quality=90)
                        img_buffer.seek(0)
                        
                        zip_file.writestr(
                            f"image_p{page_num + 1}_{img_index + 1}.jpg",
                            img_buffer.getvalue()
                        )
                        img_count += 1
                
                doc.close()
            else:
                # Convert pages to images
                images = convert_from_path(str(temp_path), dpi=dpi)
                
                for i, img in enumerate(images):
                    img_buffer = io.BytesIO()
                    img.save(img_buffer, format="JPEG", quality=90)
                    img_buffer.seek(0)
                    zip_file.writestr(f"page_{i + 1}.jpg", img_buffer.getvalue())
        
        zip_buffer.seek(0)
        return Response(
            zip_buffer.getvalue(),
            mimetype="application/zip",
            headers={"Content-Disposition": "attachment; filename=pdf_images.zip"}
        )
    finally:
        cleanup_file(temp_path)


# ==================== JPG TO PDF ====================
@app.route("/api/jpg-to-pdf", methods=["POST"])
def jpg_to_pdf():
    """Convert JPG images to PDF."""
    files = request.files.getlist("files")
    
    if not files:
        return jsonify({"detail": "No files provided"}), 400
    
    orientation = request.form.get("orientation", "portrait")
    margin = int(request.form.get("margin", 20))
    
    output_filename = generate_filename("images", "pdf")
    output_path = OUTPUT_DIR / output_filename
    temp_files = []
    
    try:
        images = []
        for file in files:
            temp_path = save_upload_file(file)
            temp_files.append(temp_path)
            
            img = Image.open(str(temp_path))
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            images.append(img)
        
        if not images:
            return jsonify({"detail": "No valid images provided"}), 400
        
        # Determine page size based on orientation
        if orientation == "landscape":
            page_size = (842, 595)  # A4 landscape
        else:
            page_size = (595, 842)  # A4 portrait
        
        # Create PDF with images
        pdf_images = []
        for img in images:
            # Resize image to fit page with margins
            max_width = page_size[0] - (2 * margin)
            max_height = page_size[1] - (2 * margin)
            
            img.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)
            
            # Create page with image centered
            page_img = Image.new("RGB", page_size, "white")
            x = (page_size[0] - img.width) // 2
            y = (page_size[1] - img.height) // 2
            page_img.paste(img, (x, y))
            pdf_images.append(page_img)
        
        # Save as PDF
        pdf_images[0].save(
            str(output_path),
            "PDF",
            resolution=100,
            save_all=True,
            append_images=pdf_images[1:] if len(pdf_images) > 1 else []
        )
        
        response = send_file(
            output_path,
            mimetype="application/pdf",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    finally:
        for temp_file in temp_files:
            cleanup_file(temp_file)


# ==================== WATERMARK ====================
@app.route("/api/watermark", methods=["POST"])
def add_watermark():
    """Add a text watermark to a PDF."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    text = request.form.get("text", "WATERMARK")
    font_size = int(request.form.get("font_size", 50))
    opacity = float(request.form.get("opacity", 0.3))
    color = request.form.get("color", "#888888")
    rotation = int(request.form.get("rotation", 45))
    position = request.form.get("position", "center")
    
    temp_path = save_upload_file(file)
    output_filename = generate_filename("watermarked", "pdf")
    output_path = OUTPUT_DIR / output_filename
    
    try:
        doc = fitz.open(str(temp_path))
        
        # Parse color
        color_hex = color.lstrip("#")
        r = int(color_hex[0:2], 16) / 255
        g = int(color_hex[2:4], 16) / 255
        b = int(color_hex[4:6], 16) / 255
        
        for page in doc:
            rect = page.rect
            
            # Calculate position
            if position == "center":
                x = rect.width / 2
                y = rect.height / 2
            elif position == "top-left":
                x = rect.width * 0.2
                y = rect.height * 0.2
            elif position == "top-right":
                x = rect.width * 0.8
                y = rect.height * 0.2
            elif position == "bottom-left":
                x = rect.width * 0.2
                y = rect.height * 0.8
            elif position == "bottom-right":
                x = rect.width * 0.8
                y = rect.height * 0.8
            else:
                x = rect.width / 2
                y = rect.height / 2
            
            # Add watermark text
            page.insert_text(
                fitz.Point(x, y),
                text,
                fontsize=font_size,
                color=(r, g, b),
                rotate=rotation,
                overlay=True
            )
        
        doc.save(str(output_path))
        doc.close()
        
        response = send_file(
            output_path,
            mimetype="application/pdf",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    finally:
        cleanup_file(temp_path)


# ==================== ROTATE PDF ====================
@app.route("/api/rotate", methods=["POST"])
def rotate_pdf():
    """Rotate PDF pages."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    rotation = int(request.form.get("rotation", 90))
    pages = request.form.get("pages", "all")
    
    temp_path = save_upload_file(file)
    output_filename = generate_filename("rotated", "pdf")
    output_path = OUTPUT_DIR / output_filename
    
    try:
        reader = PdfReader(str(temp_path))
        writer = PdfWriter()
        total_pages = len(reader.pages)
        
        # Parse page selection
        if pages == "all":
            selected_pages = set(range(total_pages))
        else:
            selected_pages = set()
            for part in pages.split(","):
                part = part.strip()
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    selected_pages.update(range(start - 1, min(end, total_pages)))
                else:
                    page_num = int(part) - 1
                    if 0 <= page_num < total_pages:
                        selected_pages.add(page_num)
        
        for i, page in enumerate(reader.pages):
            if i in selected_pages:
                page.rotate(rotation)
            writer.add_page(page)
        
        with open(output_path, "wb") as f:
            writer.write(f)
        
        response = send_file(
            output_path,
            mimetype="application/pdf",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    finally:
        cleanup_file(temp_path)


# ==================== HTML TO PDF ====================
@app.route("/api/html-to-pdf", methods=["POST"])
def html_to_pdf():
    """Convert a webpage to PDF."""
    url = request.form.get("url")
    
    if not url:
        return jsonify({"detail": "URL is required"}), 400
    
    output_filename = generate_filename("webpage", "pdf")
    output_path = OUTPUT_DIR / output_filename
    
    try:
        # Configure pdfkit options
        options = {
            'page-size': 'A4',
            'margin-top': '10mm',
            'margin-right': '10mm',
            'margin-bottom': '10mm',
            'margin-left': '10mm',
            'encoding': 'UTF-8',
            'no-outline': None,
            'enable-local-file-access': None
        }
        
        pdfkit.from_url(url, str(output_path), options=options)
        
        response = send_file(
            output_path,
            mimetype="application/pdf",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    except Exception as e:
        return jsonify({"detail": f"Conversion failed: {str(e)}"}), 500


# ==================== UNLOCK PDF ====================
@app.route("/api/unlock", methods=["POST"])
def unlock_pdf():
    """Remove password protection from a PDF."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    password = request.form.get("password", "")
    
    temp_path = save_upload_file(file)
    output_filename = generate_filename("unlocked", "pdf")
    output_path = OUTPUT_DIR / output_filename
    
    try:
        with pikepdf.open(str(temp_path), password=password) as pdf:
            pdf.save(str(output_path))
        
        response = send_file(
            output_path,
            mimetype="application/pdf",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    except pikepdf.PasswordError:
        return jsonify({"detail": "Invalid password"}), 401
    except Exception as e:
        return jsonify({"detail": f"Failed to unlock PDF: {str(e)}"}), 500
    finally:
        cleanup_file(temp_path)


# ==================== ORGANIZE PDF ====================
@app.route("/api/organize", methods=["POST"])
def organize_pdf():
    """Reorganize PDF pages - reorder, delete, or rearrange."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    page_order = request.form.get("page_order")
    delete_pages = request.form.get("delete_pages", "")
    
    if not page_order:
        return jsonify({"detail": "page_order is required"}), 400
    
    temp_path = save_upload_file(file)
    output_filename = generate_filename("organized", "pdf")
    output_path = OUTPUT_DIR / output_filename
    
    try:
        reader = PdfReader(str(temp_path))
        writer = PdfWriter()
        total_pages = len(reader.pages)
        
        # Parse pages to delete
        pages_to_delete = set()
        if delete_pages:
            for part in delete_pages.split(","):
                part = part.strip()
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    pages_to_delete.update(range(start, end + 1))
                else:
                    pages_to_delete.add(int(part))
        
        # Parse page order
        new_order = []
        for part in page_order.split(","):
            part = part.strip()
            if "-" in part:
                start, end = map(int, part.split("-"))
                new_order.extend(range(start, end + 1))
            else:
                new_order.append(int(part))
        
        # Add pages in new order, skipping deleted ones
        for page_num in new_order:
            if page_num in pages_to_delete:
                continue
            if 1 <= page_num <= total_pages:
                writer.add_page(reader.pages[page_num - 1])
        
        with open(output_path, "wb") as f:
            writer.write(f)
        
        response = send_file(
            output_path,
            mimetype="application/pdf",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    finally:
        cleanup_file(temp_path)


# ==================== PAGE NUMBERS ====================
@app.route("/api/page-numbers", methods=["POST"])
def add_page_numbers():
    """Add page numbers to a PDF."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    position = request.form.get("position", "bottom-center")
    format_str = request.form.get("format_str", "Page {n} of {total}")
    font_size = int(request.form.get("font_size", 12))
    color = request.form.get("color", "#000000")
    start_number = int(request.form.get("start_number", 1))
    skip_first = request.form.get("skip_first", "false").lower() in ("true", "1", "yes")
    
    temp_path = save_upload_file(file)
    output_filename = generate_filename("numbered", "pdf")
    output_path = OUTPUT_DIR / output_filename
    
    try:
        doc = fitz.open(str(temp_path))
        total_pages = len(doc)
        
        # Parse color
        color_hex = color.lstrip("#")
        r = int(color_hex[0:2], 16) / 255
        g = int(color_hex[2:4], 16) / 255
        b = int(color_hex[4:6], 16) / 255
        
        for i, page in enumerate(doc):
            if skip_first and i == 0:
                continue
            
            page_num = start_number + i
            if skip_first:
                page_num -= 1
            
            rect = page.rect
            margin = 30
            
            # Calculate position
            if "top" in position:
                y = margin
            else:
                y = rect.height - margin
            
            if "left" in position:
                x = margin
            elif "right" in position:
                x = rect.width - margin
            else:
                x = rect.width / 2
            
            # Format page number text
            text = format_str.replace("{n}", str(page_num)).replace("{total}", str(total_pages))
            
            page.insert_text(
                fitz.Point(x, y),
                text,
                fontsize=font_size,
                color=(r, g, b)
            )
        
        doc.save(str(output_path))
        doc.close()
        
        response = send_file(
            output_path,
            mimetype="application/pdf",
            as_attachment=True,
            download_name=output_filename
        )
        
        @response.call_on_close
        def cleanup():
            cleanup_file(output_path)
        
        return response
    finally:
        cleanup_file(temp_path)


# ==================== GET PDF INFO ====================
@app.route("/api/pdf-info", methods=["POST"])
def get_pdf_info():
    """Get information about a PDF file."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    temp_path = save_upload_file(file)
    
    try:
        doc = fitz.open(str(temp_path))
        
        info = {
            "filename": file.filename,
            "pages": len(doc),
            "metadata": doc.metadata,
            "is_encrypted": doc.is_encrypted,
            "page_sizes": []
        }
        
        for page in doc:
            rect = page.rect
            info["page_sizes"].append({
                "width": rect.width,
                "height": rect.height
            })
        
        doc.close()
        return jsonify(info)
    finally:
        cleanup_file(temp_path)


# ==================== PREVIEW / THUMBNAIL ====================
@app.route("/api/preview", methods=["POST"])
def generate_preview():
    """Generate thumbnail previews for PDF files or return image previews."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    page = int(request.form.get("page", 0))  # 0-indexed page number
    width = int(request.form.get("width", 150))  # Thumbnail width
    
    temp_path = save_upload_file(file)
    
    try:
        filename = file.filename.lower()
        
        # Handle image files
        if filename.endswith(('.jpg', '.jpeg', '.png', '.gif', '.webp')):
            img = Image.open(str(temp_path))
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            
            # Calculate height maintaining aspect ratio
            aspect = img.height / img.width
            height = int(width * aspect)
            img.thumbnail((width, height), Image.Resampling.LANCZOS)
            
            img_buffer = io.BytesIO()
            img.save(img_buffer, format="JPEG", quality=85)
            img_buffer.seek(0)
            
            import base64
            img_base64 = base64.b64encode(img_buffer.getvalue()).decode('utf-8')
            
            return jsonify({
                "thumbnail": f"data:image/jpeg;base64,{img_base64}",
                "pages": 1,
                "width": img.width,
                "height": img.height
            })
        
        # Handle PDF files
        elif filename.endswith('.pdf'):
            doc = fitz.open(str(temp_path))
            total_pages = len(doc)
            
            if page >= total_pages:
                page = 0
            
            pdf_page = doc[page]
            
            # Render page to image
            zoom = width / pdf_page.rect.width
            mat = fitz.Matrix(zoom, zoom)
            pix = pdf_page.get_pixmap(matrix=mat)
            
            img_buffer = io.BytesIO(pix.tobytes("jpeg"))
            
            import base64
            img_base64 = base64.b64encode(img_buffer.getvalue()).decode('utf-8')
            
            result = {
                "thumbnail": f"data:image/jpeg;base64,{img_base64}",
                "pages": total_pages,
                "width": pix.width,
                "height": pix.height
            }
            
            doc.close()
            return jsonify(result)
        
        else:
            return jsonify({"detail": "Unsupported file type for preview"}), 400
            
    finally:
        cleanup_file(temp_path)


@app.route("/api/preview-all-pages", methods=["POST"])
def generate_all_page_previews():
    """Generate thumbnail previews for all pages of a PDF."""
    file = request.files.get("file")
    
    if not file:
        return jsonify({"detail": "No file provided"}), 400
    
    width = int(request.form.get("width", 120))
    
    temp_path = save_upload_file(file)
    
    try:
        doc = fitz.open(str(temp_path))
        total_pages = len(doc)
        
        import base64
        thumbnails = []
        
        for i in range(total_pages):
            pdf_page = doc[i]
            
            # Render page to image
            zoom = width / pdf_page.rect.width
            mat = fitz.Matrix(zoom, zoom)
            pix = pdf_page.get_pixmap(matrix=mat)
            
            img_buffer = io.BytesIO(pix.tobytes("jpeg"))
            img_base64 = base64.b64encode(img_buffer.getvalue()).decode('utf-8')
            
            thumbnails.append({
                "page": i + 1,
                "thumbnail": f"data:image/jpeg;base64,{img_base64}",
                "width": pix.width,
                "height": pix.height
            })
        
        doc.close()
        
        return jsonify({
            "pages": total_pages,
            "thumbnails": thumbnails
        })
            
    finally:
        cleanup_file(temp_path)


# ==================== HEALTH CHECK ====================
@app.route("/api/health", methods=["GET"])
def health_check():
    """Health check endpoint."""
    return jsonify({"status": "healthy", "version": "1.0.0"})


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=8000, debug=True)
